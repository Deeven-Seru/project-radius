"""
generate_ppt.py
---------------
Generates a detailed, precise, elegant, and comprehensive PowerPoint presentation (.pptx)
for Project Radius.
Features a high-contrast dark theme (Navy Blue-Grey background, orange and ice-blue accents,
clean typography, structured tables, and native PowerPoint vector shapes/flowcharts).
No emojis are used.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# Define Color Palette
BG_COLOR = RGBColor(11, 16, 27)         # Deep Space / Navy Black (very dark)
TITLE_COLOR = RGBColor(247, 144, 8)     # Warm Orange Accent (starlight / laser)
TEXT_COLOR = RGBColor(242, 244, 247)     # Off-White
TEXT_MUTED = RGBColor(152, 162, 179)     # Light Muted Grey
ACCENT_BLUE = RGBColor(53, 185, 235)     # Ice Blue (reconstruction / DM correction)
CARD_BG = RGBColor(22, 29, 45)           # Lighter Slate Blue for diagrams / containers

def add_header(slide, title_text, subtitle_text=""):
    """Applies solid background and creates a standard slide title and subtitle."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    if subtitle_text:
        # Subtitle Text Box
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(11.833), Inches(0.4))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = 'Arial'
        p_sub.font.size = Pt(15)
        p_sub.font.color.rgb = ACCENT_BLUE

def add_two_columns(slide, left_title, left_bullets, right_title, right_bullets, bullet_height=Inches(2.3)):
    """Creates a standard two-column bulleted layout under the title area."""
    # Left Column Title
    left_title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(0.35))
    tf_lt = left_title_box.text_frame
    tf_lt.word_wrap = True
    tf_lt.margin_left = tf_lt.margin_top = tf_lt.margin_right = tf_lt.margin_bottom = 0
    p_lt = tf_lt.paragraphs[0]
    p_lt.text = left_title
    p_lt.font.name = 'Trebuchet MS'
    p_lt.font.size = Pt(18)
    p_lt.font.bold = True
    p_lt.font.color.rgb = ACCENT_BLUE
    
    # Left Column Bullets
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(5.6), bullet_height)
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
    
    for i, bullet in enumerate(left_bullets):
        p = tf_l.add_paragraph() if i > 0 else tf_l.paragraphs[0]
        p.space_after = Pt(8)
        
        if bullet.startswith("  - "):
            p.text = bullet.replace("  - ", "  • ")
            p.font.name = 'Arial'
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_MUTED
            p.level = 1
        else:
            p.text = bullet
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            
    # Right Column Title
    right_title_box = slide.shapes.add_textbox(Inches(6.983), Inches(1.6), Inches(5.6), Inches(0.35))
    tf_rt = right_title_box.text_frame
    tf_rt.word_wrap = True
    tf_rt.margin_left = tf_rt.margin_top = tf_rt.margin_right = tf_rt.margin_bottom = 0
    p_rt = tf_rt.paragraphs[0]
    p_rt.text = right_title
    p_rt.font.name = 'Trebuchet MS'
    p_rt.font.size = Pt(18)
    p_rt.font.bold = True
    p_rt.font.color.rgb = ACCENT_BLUE
    
    # Right Column Bullets
    right_box = slide.shapes.add_textbox(Inches(6.983), Inches(2.0), Inches(5.6), bullet_height)
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    
    for i, bullet in enumerate(right_bullets):
        p = tf_r.add_paragraph() if i > 0 else tf_r.paragraphs[0]
        p.space_after = Pt(8)
        
        if bullet.startswith("  - "):
            p.text = bullet.replace("  - ", "  • ")
            p.font.name = 'Arial'
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_MUTED
            p.level = 1
        else:
            p.text = bullet
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0

def add_bottom_callout(slide, title_text, description_text):
    """Draws a callout box at the bottom of the slide to summarize key takeaway."""
    left = Inches(0.75)
    top = Inches(5.8)
    width = Inches(11.833)
    height = Inches(1.1)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = TITLE_COLOR
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.12)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_bottom = Inches(0.12)
    
    p = tf.paragraphs[0]
    p.text = title_text + ": "
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Add inline description text to same paragraph
    run = p.add_run()
    run.text = description_text
    run.font.name = 'Arial'
    run.font.size = Pt(13)
    run.font.color.rgb = TEXT_COLOR
    run.font.bold = False

def add_architecture_diagram(slide):
    """Draws a 4-step flowchart showing the data processing pipeline."""
    steps = [
        ("Acquisition HAL", "GenTL / RAM Cache\nZero-copy buffer access"),
        ("slopes.c (C)", "Center of Gravity (CoG)\n632 X/Y Slope Vectors"),
        ("reconstruction.c (C)", "Matrix Multiplication\n55 Zernike Modes"),
        ("dm_coupling.csv (C)", "Inter-actuator Matrix\n357 DM Actuator Volts")
    ]
    
    start_left = Inches(0.75)
    top = Inches(4.3)
    box_width = Inches(2.2)
    box_height = Inches(1.3)
    arrow_width = Inches(0.5)
    arrow_height = Inches(0.3)
    gap = Inches(0.75)
    
    for i, (title, desc) in enumerate(steps):
        left = start_left + i * (box_width + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = ACCENT_BLUE
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.08)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_COLOR
        p2.space_before = Pt(5)
        p2.alignment = PP_ALIGN.CENTER
        
        if i < len(steps) - 1:
            arrow_left = left + box_width + gap/2 - arrow_width/2
            arrow_top = top + box_height/2 - arrow_height/2
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, arrow_width, arrow_height)
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TITLE_COLOR
            arrow.line.fill.background()

def add_hal_diagram(slide):
    """Draws 3 side-by-side cards representing camera interfaces."""
    modes = [
        ("Simulated Mode (1000 Hz)", "Runs background physical optics simulation inside camera_interface.py. Populates dummy frame buffers via light propagation equations."),
        ("Playback Mode (Disk-Free)", "Pre-loads optical dataset frames directly into RAM. Feeds buffers in a thread loop, removing disk I/O bottlenecks to test C-Engine limits."),
        ("Hardware Mode (GenICam)", "Uses EMVA standard harvesters interface to link live Camera Link or CoaXPress hardware using manufacturer-provided .cti drivers.")
    ]
    
    start_left = Inches(0.75)
    top = Inches(4.3)
    box_width = Inches(3.4)
    box_height = Inches(1.3)
    gap = Inches(0.8)
    
    for i, (title, desc) in enumerate(modes):
        left = start_left + i * (box_width + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = ACCENT_BLUE
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.08)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_COLOR
        p2.space_before = Pt(5)
        p2.alignment = PP_ALIGN.CENTER

def add_math_flow_diagram(slide):
    """Draws math transformation blocks showing the physics flow."""
    steps = [
        ("Wavefront Phase Phi", "Distorted incoming wavefront\n(Fried parameter r0)"),
        ("Subaperture Gradients (S)", "MLA subdivisions: 316 local\ntilt displacements in X & Y"),
        ("Zernike Modes (a)", "Expansion over circular pupil:\n55 modes (Z2 to Z55)"),
        ("Actuator Voltages (V)", "MVM reconstruction: V = G+ * S\nmapping to 357 DM actuators")
    ]
    
    start_left = Inches(0.75)
    top = Inches(4.3)
    box_width = Inches(2.2)
    box_height = Inches(1.3)
    arrow_width = Inches(0.5)
    arrow_height = Inches(0.3)
    gap = Inches(0.75)
    
    for i, (title, desc) in enumerate(steps):
        left = start_left + i * (box_width + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = ACCENT_BLUE
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.08)
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_COLOR
        p2.space_before = Pt(5)
        p2.alignment = PP_ALIGN.CENTER
        
        if i < len(steps) - 1:
            arrow_left = left + box_width + gap/2 - arrow_width/2
            arrow_top = top + box_height/2 - arrow_height/2
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, arrow_width, arrow_height)
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TITLE_COLOR
            arrow.line.fill.background()

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Elegant design)
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Decorative colored accent bars
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(4.5), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TITLE_COLOR
    accent_bar.line.fill.background()
    
    accent_bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.35), Inches(2.0), Inches(0.04))
    accent_bar2.fill.solid()
    accent_bar2.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar2.line.fill.background()
    
    # Title text box
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.6), Inches(11.833), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "PROJECT RADIUS"
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    p_sub = tf.add_paragraph()
    p_sub.text = "High-Performance C-Engine for Real-Time Wavefront Reconstruction & Turbulence Characterization"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = ACCENT_BLUE
    p_sub.space_before = Pt(8)
    
    p_auth = tf.add_paragraph()
    p_auth.text = "System Architect: Deeven Seru  |  Widescreen POSIX-compliant C/Python Integration"
    p_auth.font.name = 'Arial'
    p_auth.font.size = Pt(14)
    p_auth.font.color.rgb = TEXT_MUTED
    p_auth.space_before = Pt(36)
    
    # -------------------------------------------------------------
    # SLIDE 2: The Atmospheric Challenge
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "The Atmospheric Challenge", "Physical optical disturbances and hardware systems correction limits")
    add_two_columns(
        slide,
        "Physical Optical Turbulence",
        [
            "Index of refraction fluctuations in atmospheric layers distort plane-parallel light wavefronts.",
            "Causes severe beam wander, focal blur, and communication phase errors.",
            "Fried Parameter (r0): Defines transverse coherence diameter of the incoming wavefront (typically 5-20 cm).",
            "Coherence Time (tau0): Characteristic timeframe over which turbulence structures evolve (typically 2-10 ms)."
        ],
        "Systems Engineering Bottlenecks",
        [
            "The 10 ms Real-Time Deadline: Command loop must run faster than tau0 to achieve high Strehl corrections.",
            "High-speed Shack-Hartmann Sensors (WFS): Capture frames at 1000 Hz with minimal readout overhead.",
            "Software OS Latency: Standard operating systems and high-level languages introduce non-deterministic garbage collection delays.",
            "Real-Time Constraint: C-Engine is required to perform deterministic sub-millisecond execution."
        ]
    )
    add_bottom_callout(
        slide,
        "Physical Loop Constraint",
        "The coherence time tau0 sets a strict 10 ms upper boundary. Loop latency must remain sub-millisecond (<1.0 ms) to prevent control lag and correct higher-order aberrations."
    )
    
    # -------------------------------------------------------------
    # SLIDE 3: Scientific Foundations
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Scientific Foundations", "First-principles physical modeling of wavefront sensing and deformable mirror mapping")
    add_two_columns(
        slide,
        "Wavefront Sensing (SH-WFS)",
        [
            "Microlens Array (MLA) divides the incoming pupil wavefront into individual focal spots.",
            "Local wavefront tilt causes spot displacements on the camera sensor pixel array.",
            "Centroid coordinate offsets calculate the wavefront gradient.",
            "Gradient Vector (S): 632 coordinates (316 valid circular aperture subapertures in X and Y directions)."
        ],
        "Reconstruction & Actuation",
        [
            "Zernike Polynomials (Noll 1976): Circular-pupil orthogonal modes mapping physical aberrations (Z2 to Z55).",
            "Control Matrix G+: Relates slope vectors to Zernike modes and DM actuator commands via Least-Squares regression.",
            "Deformable Mirror (DM): 357 actuators adjust facesheet shape using inter-actuator coupling coefficients."
        ]
    )
    add_math_flow_diagram(slide)
    
    # -------------------------------------------------------------
    # SLIDE 4: Software Architecture (C-Engine)
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Software Architecture: The C-Engine", "Bare-metal latency and zero-copy pipeline integration")
    add_two_columns(
        slide,
        "Optimized C Execution Layer",
        [
            "Centroiding and Matrix-Vector Multiplication (MVM) are implemented in pure POSIX C (c_engine.so).",
            "Pre-allocated static structs (geometry.h) eliminate real-time heap allocation overhead.",
            "Compilation optimized for AVX2 vector registers to process centroids and MVM in parallel.",
            "Total pipeline execution time: 0.31 ms (enabling >3000 Hz loop frequencies)."
        ],
        "Zero-Copy Ctypes Bridge",
        [
            "Python handles non-real-time orchestration, file IO, and statistics compilation.",
            "ctypes bridge passes raw camera buffer memory addresses directly to C structures.",
            "Avoids intermediate data copies, maintaining maximum CPU cache locality.",
            "Ensures deterministic processing, avoiding standard operating system thread scheduling latency."
        ]
    )
    add_architecture_diagram(slide)
    
    # -------------------------------------------------------------
    # SLIDE 5: Live Camera & HAL
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Live Camera & HAL Integration", "Hardware Abstraction Layer for lab simulations and observatory cameras")
    add_two_columns(
        slide,
        "Acquisition HAL (camera_interface.py)",
        [
            "Abstracts frame grabber operations behind a standardized, pluggable class interface.",
            "Simulated Mode: Generates synthetic Shack-Hartmann spot arrays using physical optics math.",
            "Playback Mode: Pre-loads experimental datasets directly into RAM for high-speed playback.",
            "Hardware Mode: Integrates with physical GenICam / Camera Link cameras."
        ],
        "EMVA GenICam & GenTL Link",
        [
            "Uses the Harvesters framework to load manufacturer .cti (GenTL) driver libraries.",
            "Supports CoaXPress, GigE Vision, USB3 Vision, and Camera Link frame grabbers.",
            "Bypasses standard python buffer allocations by copying DMA memory straight to the C-Engine.",
            "Compatible with professional cameras from Basler, Euresys, Active Silicon, and FLIR."
        ]
    )
    add_hal_diagram(slide)
    
    # -------------------------------------------------------------
    # SLIDE 6: Results & Benchmarks
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Results & Achievements (55 Zernike Modes)", "Performance metrics under 1000 Hz playback loop")
    
    # Add table for results
    rows, cols = 7, 4
    left, top, width, height = Inches(0.75), Inches(1.6), Inches(11.833), Inches(4.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set Column Widths
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(2.833)
    
    # Define Data
    headers = ["Metric", "Target Requirement", "Achieved (55 Modes)", "Status"]
    data = [
        ["End-to-End Processing Latency", "< 10.00 ms", "0.309 ms", "Pass"],
        ["Average Closed-Loop Rate", "1000 Hz", "998.05 Hz", "Pass"],
        ["Average Per-Frame Spatial Accuracy (R²)", "> 95.00%", "99.4857%", "Pass"],
        ["Absolute Global Temporal Accuracy (R²)", "> 95.00%", "98.6009%", "Pass"],
        ["Reconstruction Mean Square Error (MSE)", "< 0.1", "1.14e-15", "Pass"],
        ["Reconstructed Strehl Ratio", "Maximize", "93.26% – 96.76%", "Pass"]
    ]
    
    # Style Headers
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = BG_COLOR
        
    # Style Data Cells
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            p.font.name = 'Arial'
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_COLOR
            if col_idx == 3:
                p.font.bold = True
                p.font.color.rgb = ACCENT_BLUE
                
    # Add a small note at the bottom
    note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.25), Inches(11.833), Inches(0.4))
    tf_note = note_box.text_frame
    tf_note.word_wrap = True
    p_note = tf_note.paragraphs[0]
    p_note.text = "Note: Tests conducted on a multi-threaded 1000 Hz frame stream using pre-loaded experimental calibration configurations."
    p_note.font.name = 'Arial'
    p_note.font.size = Pt(11)
    p_note.font.italic = True
    p_note.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 7: Unbiased Robustness Analysis
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Unbiased Robustness Analysis", "Quantifying closed-loop performance degradation under physical noise and drift")
    add_two_columns(
        slide,
        "Noise Sensitivities (Photon vs Readout)",
        [
            "Photon Noise: Highly Robust. The modal Zernike reconstructor acts as a spatial low-pass filter, preserving 99.14% R² at 200 photons/subaperture/frame.",
            "Readout Noise (RON): High Sensitivity. Background pixel noise shifts Center of Gravity (CoG) measurements toward the subaperture center.",
            "  - At 1000 photons/subap, R² degrades to 97.80% under 1.0 e- RON (sCMOS) and drops to 75.29% under 5.0 e- RON (CCD)."
        ],
        "Alignment Drift Sensitivity",
        [
            "Lenslet Spot Field Drift: Extreme Sensitivity.",
            "A uniform alignment shift of just 0.05 pixels reduces R² to 59.65%; shifts of >=0.10 pixels collapse the loop entirely.",
            "  - Uniform spot shifts are mathematically reconstructed as a large artificial Tip/Tilt error, driving the DM actuators to their physical limits.",
            "Mitigation: Requires real-time reference slope updates and active Tip/Tilt mirror path stabilization."
        ]
    )
    add_bottom_callout(
        slide,
        "Robustness Summary",
        "While photon noise is naturally smoothed by Zernike reconstruction, high readout noise and subpixel spot drift are critical loop failure vectors that require active background thresholding and tip/tilt alignment stabilization."
    )

    # -------------------------------------------------------------
    # SLIDE 8: Cost-Benefit & Scalability Analysis
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Space Agency (ISRO) Cost-Benefit & Scalability", "Analyzing economic and technical value for satellite ground stations and observatories")
    add_two_columns(
        slide,
        "Drastic Hardware Cost Reductions",
        [
            "Proprietary RTC Hardware Bypass: Eliminates the need for expensive custom FPGA/ASIC real-time controllers (saving $100k+ per observatory installation).",
            "Commodity COTS Servers: Runs on standard x86_64 or ARM64 Linux servers. The sub-millisecond C-Engine execution leaves CPU capacity free for other systems.",
            "Open Source Ecosystem: Built using free, modular components (POSIX, Python, ctypes), eliminating annual software licensing fees."
        ],
        "Scalability & Standards Compliance",
        [
            "Zernike Mode Scaling: Successfully validated scaling from 20 to 55 modes, capturing high-order turbulence details without latency penalties.",
            "Universal Camera Compatibility: Full compliance with the EMVA GenICam standard ensures interoperability with any CoaXPress/GigE camera.",
            "Ground Station Integration: Portable C-library structure simplifies integration into satellite communications and tracking controllers."
        ]
    )
    add_bottom_callout(
        slide,
        "Strategic Value",
        "Project Radius proves that high-order (55-mode) adaptive optics RTC pipelines can run on standard POSIX systems. This slashes developmental and hardware costs while ensuring scalability to larger systems."
    )

    # -------------------------------------------------------------
    # SLIDE 9: Future Work & Conclusions
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Future Work & Conclusions", "Scaling adaptive optics reconstruction to the next generation")
    add_two_columns(
        slide,
        "Technical Development Roadmap",
        [
            "Closed-Loop Control Integration: Implement PI (Proportional-Integral) loop filters to dynamically command Deformable Mirrors.",
            "Dynamic Centroid Thresholding: Add adaptive background subtraction to centroid computations to completely reject read noise (RON).",
            "OpenMP Multi-Core Scaling: Parallelize the Matrix-Vector Multiplication (MVM) loop to support high-density subaperture grids (e.g. 80x80)."
        ],
        "Key Project Takeaways",
        [
            "Elite Performance: 0.31 ms loop latency matches dedicated hardware real-time controller performance.",
            "High Accuracy: 99.49% spatial R² and 98.60% temporal accuracy confirm the precision of the 55-Zernike-mode engine.",
            "Production Ready: Modular architecture ready for laboratory and observatory deployments."
        ]
    )
    add_bottom_callout(
        slide,
        "Conclusion",
        "Project Radius provides a high-performance, open-standard, low-latency foundation for real-time wavefront reconstruction, suitable for space agency ground stations and astronomical observatories."
    )
    
    prs.save(os.path.join(BASE, 'project_radius_presentation.pptx'))
    print("PowerPoint presentation generated successfully: project_radius_presentation.pptx")

if __name__ == '__main__':
    main()
